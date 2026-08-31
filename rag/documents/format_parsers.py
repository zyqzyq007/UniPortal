"""
Multi-format document parsers: DOCX, PPTX, HTML, plus plain-text fallback.

Each parser returns ``List[Document]`` (one Document per paragraph / slide /
section) with ``source`` metadata set to the filename. Parsers are
dependency-soft: the required library (python-docx / python-pptx /
beautifulsoup4) is imported lazily so the module loads even when the library
is absent — a clear ``RuntimeError`` is raised only when that format is
actually requested.

This keeps the offline/default deployment free of extra heavy deps while
making the formats available wherever the libs are installed.
"""

from __future__ import annotations

from langchain_core.documents import Document

from utils.log_utils import log

__all__ = ["parse_docx", "parse_pptx", "parse_html", "parse_by_extension"]


def parse_docx(file_path: str, source: str = "") -> list[Document]:
    """Parse a .docx file into paragraph-level documents."""
    try:
        from docx import Document as _DocxDocument  # python-docx
    except ImportError as e:
        raise RuntimeError("DOCX parsing requires python-docx: pip install python-docx") from e

    doc = _DocxDocument(file_path)
    docs: list[Document] = []
    for para in doc.paragraphs:
        text = (para.text or "").strip()
        if len(text) < 4:
            continue
        style = para.style.name if para.style else ""
        docs.append(
            Document(
                page_content=text,
                metadata={"source": source, "section": style, "format": "docx"},
            )
        )
    log.info(f"DOCX parsed: {len(docs)} paragraphs from {source or file_path}")
    return docs


def parse_pptx(file_path: str, source: str = "") -> list[Document]:
    """Parse a .pptx file into slide-level documents."""
    try:
        from pptx import Presentation  # python-pptx
    except ImportError as e:
        raise RuntimeError("PPTX parsing requires python-pptx: pip install python-pptx") from e

    prs = Presentation(file_path)
    docs: list[Document] = []
    for idx, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(run.text for run in para.runs).strip()
                    if t:
                        texts.append(t)
            elif getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        texts.append(" | ".join(cells))
        body = "\n".join(texts).strip()
        if len(body) < 4:
            continue
        docs.append(
            Document(
                page_content=body,
                metadata={
                    "source": source,
                    "page": idx,
                    "section": f"幻灯片{idx}",
                    "format": "pptx",
                },
            )
        )
    log.info(f"PPTX parsed: {len(docs)} slides from {source or file_path}")
    return docs


def parse_html(file_path: str, source: str = "") -> list[Document]:
    """Parse an HTML file into section-level documents."""
    try:
        from bs4 import BeautifulSoup  # beautifulsoup4
    except ImportError as e:
        raise RuntimeError(
            "HTML parsing requires beautifulsoup4: pip install beautifulsoup4"
        ) from e

    with open(file_path, encoding="utf-8", errors="ignore") as f:
        html = f.read()
    soup = BeautifulSoup(html, "html.parser")
    # Remove script/style noise.
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()

    docs: list[Document] = []
    # Split by headings to preserve structure.
    current_title = ""
    buffer: list[str] = []
    headings = {"h1", "h2", "h3", "h4"}

    def _flush():
        nonlocal buffer
        text = "\n".join(buffer).strip()
        if len(text) >= 4:
            docs.append(
                Document(
                    page_content=text,
                    metadata={
                        "source": source,
                        "title": current_title,
                        "format": "html",
                    },
                )
            )
        buffer = []

    for el in soup.find_all(["h1", "h2", "h3", "h4", "p", "li", "table"]):
        if el.name in headings:
            _flush()
            current_title = el.get_text(strip=True)
        else:
            t = el.get_text(" ", strip=True)
            if t:
                buffer.append(t)
    _flush()

    log.info(f"HTML parsed: {len(docs)} sections from {source or file_path}")
    return docs


def parse_by_extension(file_path: str, source: str = "") -> list[Document]:
    """
    Dispatch to the right parser by file extension.

    Raises ValueError for unknown extensions.
    """
    import os

    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".docx":
        return parse_docx(file_path, source=source)
    if ext == ".pptx":
        return parse_pptx(file_path, source=source)
    if ext in (".html", ".htm"):
        return parse_html(file_path, source=source)
    raise ValueError(f"Unsupported extension for multi-format parser: {ext}")
